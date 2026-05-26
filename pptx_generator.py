from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import json

def generate_pptx_stream(company_name, stream):
    from ai_engine import _call_gemini
    
    prompt = f"""
    Eres un consultor de ventas de tecnología experto en IA y automatización. Necesito estructurar una presentación de ventas de 6 diapositivas para convencer a la clínica "{company_name}" de implementar un asistente de Inteligencia Artificial de Egasis Studio para su atención en WhatsApp y agendamiento.
    
    Genera el contenido estructurado en formato JSON válido.
    Debe tener exactamente esta estructura de claves:
    {{
        "slide1_title": "...",
        "slide1_subtitle": "...",
        "slide2_title": "...",
        "slide2_bullets": ["...", "...", "..."],
        "slide3_title": "...",
        "slide3_bullets": ["...", "...", "..."],
        "slide4_title": "...",
        "slide4_bullets": ["...", "...", "..."],
        "slide5_title": "...",
        "slide5_bullets": ["...", "...", "..."],
        "slide6_title": "...",
        "slide6_bullets": ["...", "...", "..."]
    }}
    
    Por favor, sé muy específico con el sector dental/médico y con la clínica "{company_name}". 
    Diapositivas:
    1. Portada (Título y subtítulo llamativo y profesional sobre transformación digital)
    2. El Diagnóstico (Pérdida de pacientes en WhatsApp, respuestas fuera de hora, staff dental colapsado)
    3. La Solución (Asistente de IA a medida: responde 24/7 y agenda solo)
    4. Beneficios Clave (Aumento de reservas de tratamientos estéticos, ahorro de tiempo del staff, retención 100%)
    5. Plan de Acción (Fase 1: Configuración; Fase 2: Integración; Fase 3: Lanzamiento)
    6. Siguiente Paso (Llamado a la acción para iniciar el despliegue en 2 semanas)
    
    Responde ÚNICAMENTE con el bloque JSON crudo. No incluyas explicaciones, ni ```json ni ```.
    """
    
    content = None
    try:
        raw_res = _call_gemini(prompt)
        # Limpiar markdown
        raw_res = raw_res.strip()
        if raw_res.startswith("```json"):
            raw_res = raw_res[7:]
        elif raw_res.startswith("```"):
            raw_res = raw_res[3:]
        if raw_res.endswith("```"):
            raw_res = raw_res[:-3]
        raw_res = raw_res.strip()
        content = json.loads(raw_res)
    except Exception as e:
        print("Error al generar contenido PPTX con Gemini, usando fallback:", e)
        
    # Fallback si falla el JSON
    if not content or not isinstance(content, dict):
        content = {
            "slide1_title": "Transformación Digital e Inteligencia Artificial",
            "slide1_subtitle": f"Propuesta exclusiva de automatización para {company_name}\nDesarrollada por Egasis Studio",
            "slide2_title": "El Desafío Actual en la Atención",
            "slide2_bullets": [
                "Pérdida de potenciales pacientes por demoras en respuestas en WhatsApp.",
                "Falta de atención y agendamiento activo fuera del horario de la clínica.",
                "El personal administrativo consume horas gestionando reprogramaciones manuales."
            ],
            "slide3_title": "La Solución: Asistente Autónomo Egasis",
            "slide3_bullets": [
                "Atención automatizada por WhatsApp activa las 24 horas, los 7 días de la semana.",
                "Cualificación instantánea del interés de pacientes en tratamientos estéticos y generales.",
                "Sincronización automática de turnos en la agenda médica sin intervención humana."
            ],
            "slide4_title": "Beneficios de la Automatización",
            "slide4_bullets": [
                "Incremento estimado del 25% en reserva de turnos de primera consulta.",
                "Liberación del 80% de tareas repetitivas en recepción.",
                "Respuestas instantáneas y personalizadas, elevando la experiencia del paciente."
            ],
            "slide5_title": "Plan de Implementación Rápido",
            "slide5_bullets": [
                "Fase 1: Configuración de la base de conocimientos y diálogos del asistente.",
                "Fase 2: Pruebas y validación en entorno privado y sincronización de agenda.",
                "Fase 3: Integración oficial en la línea de WhatsApp de la clínica."
            ],
            "slide6_title": "Siguientes Pasos",
            "slide6_bullets": [
                "Validación de tratamientos prioritarios y agenda de turnos demo.",
                "Tiempo estimado de implementación: 2 semanas de desarrollo.",
                "Comenzamos la transformación de la atención para liderar en el sector."
            ]
        }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Colores de la paleta dark premium
    BG_COLOR = RGBColor(10, 15, 29) # Slate/Navy profundo
    TEAL_COLOR = RGBColor(45, 212, 191) # Tech Teal
    WHITE_COLOR = RGBColor(255, 255, 255)
    
    # ------------------ SLIDE 1: PORTADA ------------------
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = BG_COLOR
    
    # Decoración: una franja lateral
    side_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = TEAL_COLOR
    side_bar.line.fill.background()
    
    # Título Portada
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.get("slide1_title", "Transformación Digital")
    p.font.name = 'Helvetica'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR
    
    # Subtítulo Portada
    p2 = tf.add_paragraph()
    p2.text = content.get("slide1_subtitle", f"Propuesta exclusiva para {company_name}")
    p2.font.name = 'Helvetica'
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEAL_COLOR
    p2.space_before = Pt(20)
    
    # ------------------ SLIDES DE CONTENIDO (2 a 6) ------------------
    slides_data = [
        {"title_key": "slide2_title", "bullets_key": "slide2_bullets"},
        {"title_key": "slide3_title", "bullets_key": "slide3_bullets"},
        {"title_key": "slide4_title", "bullets_key": "slide4_bullets"},
        {"title_key": "slide5_title", "bullets_key": "slide5_bullets"},
        {"title_key": "slide6_title", "bullets_key": "slide6_bullets"},
    ]
    
    for idx, slide_info in enumerate(slides_data, start=2):
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_COLOR
        
        # Barra decorativa inferior
        footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = TEAL_COLOR
        footer_bar.line.fill.background()
        
        # Título diapositiva
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.get(slide_info["title_key"], f"Slide {idx}")
        p.font.name = 'Helvetica'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEAL_COLOR
        
        # Viñetas
        bullets = content.get(slide_info["bullets_key"], [])
        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        for bullet_idx, bullet_text in enumerate(bullets):
            p_bullet = tf_content.add_paragraph() if bullet_idx > 0 else tf_content.paragraphs[0]
            p_bullet.text = bullet_text
            p_bullet.font.name = 'Helvetica'
            p_bullet.font.size = Pt(18)
            p_bullet.font.color.rgb = WHITE_COLOR
            p_bullet.space_after = Pt(16)
            p_bullet.level = 0
            
    prs.save(stream)

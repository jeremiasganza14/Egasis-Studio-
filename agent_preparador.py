import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import google.generativeai as genai
from dotenv import load_dotenv

load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY", ""))

def get_hq_path():
    desktop_hq = "/Users/jereganza/Desktop/Egasis_HQ"
    try:
        os.makedirs(desktop_hq, exist_ok=True)
    except Exception:
        desktop_hq = os.path.abspath(os.path.join(os.path.dirname(__file__), "Egasis_HQ"))
        os.makedirs(desktop_hq, exist_ok=True)
    return desktop_hq

HQ_PATH = get_hq_path()
LEADS_DIR = os.path.join(HQ_PATH, "1_Leads")
REUNIONES_DIR = os.path.join(HQ_PATH, "2_Reuniones")
PPTX_DIR = os.path.join(HQ_PATH, "3_Presentaciones")

def create_pptx(company_name, output_path):
    prs = Presentation()
    
    # 1. Slide Título
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Automatización con Inteligencia Artificial"
    subtitle.text = f"Propuesta exclusiva para {company_name}\nDesarrollada por Egasis Studio"
    
    # 2. Slide El Problema
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "¿Cuál es el desafío de hoy?"
    tf = body_shape.text_frame
    tf.text = "Las clínicas pierden pacientes por:"
    p = tf.add_paragraph()
    p.text = "Respuestas lentas en WhatsApp."
    p = tf.add_paragraph()
    p.text = "Atención al cliente limitada fuera del horario laboral."
    p = tf.add_paragraph()
    p.text = "Gestión manual de agenda que quita tiempo al staff."
    
    # 3. Slide La Solución
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "La Solución: IA Autónoma"
    tf = body_shape.text_frame
    tf.text = "Sistemas Privados implementados en sus servidores:"
    p = tf.add_paragraph()
    p.text = "Responde como un humano 24/7."
    p = tf.add_paragraph()
    p.text = "Agenda reuniones directamente en Google Calendar/Calendly."
    p = tf.add_paragraph()
    p.text = "Califica al lead antes de mandarlo a la recepcionista."
    
    # 4. Slide Siguiente Paso
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "¿Comenzamos la integración?"
    subtitle.text = "Tiempo estimado de despliegue: 2 semanas."
    
    prs.save(output_path)

def generate_call_brief(company_name, website, output_path):
    prompt = f"""
Eres un asesor de ventas B2B experto en automatización para clínicas.
Tengo una reunión con la empresa: {company_name}.
Su web es: {website}

Redacta un "Call Brief" en Markdown con la siguiente estructura:
## 1. Contexto de la Empresa
(Adivina o resume qué hacen basándote en el nombre y web).
## 2. Pros (Puntos fuertes)
## 3. Contras (Dolores probables que la IA resolverá)
## 4. Ángulo de Venta
(Cómo debo venderles nuestro sistema de IA)
    """
    
    try:
        model = genai.GenerativeModel('gemini-2.5-flash')
        response = model.generate_content(prompt)
        content = response.text
    except Exception as e:
        content = f"Error conectando a Gemini: {e}\n\nPrepara la venta enfocándote en la automatización de WhatsApp."
        
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(f"# Call Brief: {company_name}\n\n")
        f.write(content)

def check_agendados():
    print("👔 [Agente Preparador] Buscando leads agendados...")
    for filepath in glob.glob(os.path.join(LEADS_DIR, "*.md")):
        with open(filepath, "r", encoding="utf-8") as f:
            content = f.read()
            
        if "status: agendado" in content:
            # Extraer info
            company = ""
            website = ""
            for line in content.splitlines():
                if line.startswith("company:"): company = line.replace("company:", "").strip()
                if line.startswith("website:"): website = line.replace("website:", "").strip()
                
            print(f"   => Preparando material para: {company}")
            
            # 1. Crear PPTX
            pptx_filename = f"Pitch_{company.replace(' ', '_')}.pptx"
            pptx_path = os.path.join(PPTX_DIR, pptx_filename)
            create_pptx(company, pptx_path)
            print(f"      ✅ PPTX generado: {pptx_filename}")
            
            # 2. Crear Call Brief
            brief_filename = f"Brief_{company.replace(' ', '_')}.md"
            brief_path = os.path.join(REUNIONES_DIR, brief_filename)
            generate_call_brief(company, website, brief_path)
            print(f"      ✅ Call Brief generado: {brief_filename}")
            
            # 3. Actualizar status a "preparado"
            content = content.replace("status: agendado", "status: preparado")
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)

if __name__ == "__main__":
    check_agendados()
